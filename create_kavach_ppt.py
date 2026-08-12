#!/usr/bin/env python3
"""Generate KAVACH AI Presentation as PowerPoint — same content as HTML version."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════════════════════
# COLORS & HELPERS
# ═══════════════════════════════════════════════════════════
BG = RGBColor(0x0A, 0x0E, 0x1A)
PURPLE = RGBColor(0x81, 0x8C, 0xF8)
LIGHT_PURPLE = RGBColor(0xA5, 0xB4, 0xFC)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
DIM = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xF8, 0x71, 0x71)
GREEN = RGBColor(0x34, 0xD3, 0x99)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
CARD_BG = RGBColor(0x16, 0x1B, 0x30)
BORDER = RGBColor(0x33, 0x41, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def gradient_bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Pt(4))
    s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()

def card(slide, l, t, w, h, border_color=BORDER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = CARD_BG
    s.line.color.rgb = border_color; s.line.width = Pt(1.5)
    return s

def txt(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Segoe UI"
    p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines, size=14, color=GRAY, spacing=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (text, c, b, s) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(s if s else size)
        p.font.color.rgb = c if c else color
        p.font.bold = b; p.font.name = "Segoe UI"
        p.space_after = Pt(spacing)

def slide_num(slide, num):
    txt(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
        f"{num:02d} / 10", 11, False, DIM, PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(1), Inches(0.8), Inches(8), Inches(0.4), "Hexaware Profuellers  |  Agentic Arena 2026", 12, False, PURPLE)
txt(s, Inches(1), Inches(2.2), Inches(10), Inches(1), "KAVACH AI", 60, True, WHITE)
txt(s, Inches(1), Inches(3.5), Inches(11), Inches(0.7), "Kontinuous Audit & Vulnerability Analysis for Compliant Hardening", 24, False, LIGHT_PURPLE)
txt(s, Inches(1), Inches(4.5), Inches(10), Inches(0.5), '"One commit. Six domains. Zero breaches."', 16, False, DIM)
# Tags
multi(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5), [
    ("Multi-Agent AI   |   Custom Fine-Tuned LLM   |   Graph RAG   |   Real-Time Compliance", GRAY, False, 13),
], size=13)
slide_num(s, 1)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "THE PROBLEM", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.8), "One Code Commit Can Simultaneously Violate 6 Compliance Domains", 32, True, WHITE)

# Card 1: Today's Reality
card(s, Inches(0.8), Inches(2.2), Inches(3.8), Inches(4.5), RED)
multi(s, Inches(1.1), Inches(2.4), Inches(3.4), Inches(4.2), [
    ("Today's Reality", RED, True, 18),
    ("", GRAY, False, 6),
    ("SOX team finds issues in annual audit", GRAY, False, 13),
    ("Security team scans weekly", GRAY, False, 13),
    ("Legal reviews contracts quarterly", GRAY, False, 13),
    ("Fair lending tested semi-annually", GRAY, False, 13),
    ("", GRAY, False, 6),
    ("Violations found MONTHS later", RED, True, 14),
])

# Card 2: The Gap
card(s, Inches(4.9), Inches(2.2), Inches(3.8), Inches(4.5), AMBER)
multi(s, Inches(5.2), Inches(2.4), Inches(3.4), Inches(4.2), [
    ("The Gap", AMBER, True, 18),
    ("", GRAY, False, 6),
    ("Different teams, different tools", GRAY, False, 13),
    ("No cross-domain reasoning", GRAY, False, 13),
    ("No per-client context", GRAY, False, 13),
    ("Reactive, not predictive", GRAY, False, 13),
    ("", GRAY, False, 6),
    ("$50K+ per missed violation", AMBER, True, 14),
])

# Card 3: Solution
card(s, Inches(9.0), Inches(2.2), Inches(3.8), Inches(4.5), GREEN)
multi(s, Inches(9.3), Inches(2.4), Inches(3.4), Inches(4.2), [
    ("KAVACH AI Solution", GREEN, True, 18),
    ("", GRAY, False, 6),
    ("AI reasons across ALL domains simultaneously", GRAY, False, 13),
    ("Per-engagement compliance twin", GRAY, False, 13),
    ("Predicts impact BEFORE deploy", GRAY, False, 13),
    ("Generates audit evidence in real-time", GRAY, False, 13),
    ("", GRAY, False, 6),
    ("Zero compliance surprises", GREEN, True, 14),
])
slide_num(s, 2)

# ═══════════════════════════════════════════════════════════
# SLIDE 3: Architecture
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "PLATFORM ARCHITECTURE", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "End-to-End Agentic AI Pipeline", 32, True, WHITE)
txt(s, Inches(0.8), Inches(1.5), Inches(12), Inches(0.4), "Every event flows: Perceive -> Reason -> Decide -> Act. No human in the loop.", 13, False, GRAY)

# Architecture flow boxes
flow_y = [2.1, 2.9, 3.7, 4.9, 5.7]
flow_data = [
    ("PERCEIVE:  Event Sources", "Git  |  Jenkins  |  Jira  |  AWS Config  |  File Watcher  |  SAST Tools", BORDER, PURPLE),
    ("ActiveMQ Message Queue - \"compliance-events\"", "Unified event bus: all sources flow into a single queue", PURPLE, LIGHT_PURPLE),
    ("REASON:  Agent Layer", "Digital Twin  <-  Chain Reactor (central)  ->  Drift Sentinel  |  Audit Narrator  <->  Obligation Parser", GREEN, GREEN),
    ("DECIDE:  Intelligence Layer", "KAVACH LLM (our model)  |  Knowledge Graph + GraphRAG  |  Policy Engine (26+ rules)", LIGHT_PURPLE, LIGHT_PURPLE),
    ("ACT:  Actions & Outputs", "Deploy Gate (Block/Allow)  |  Audit Evidence (auto)  |  Dashboard  |  Alerts (Slack/Email)", RED, RED),
]
for i, (title, sub, border, title_color) in enumerate(flow_data):
    y = Inches(flow_y[i])
    card(s, Inches(1.5), y, Inches(10.3), Inches(0.7), border)
    txt(s, Inches(1.8), y + Pt(6), Inches(9.5), Inches(0.35), title, 14, True, title_color)
    txt(s, Inches(1.8), y + Pt(30), Inches(9.5), Inches(0.3), sub, 11, False, GRAY)
    # Arrow between boxes
    if i < len(flow_data) - 1:
        txt(s, Inches(6.5), y + Inches(0.7), Inches(0.5), Inches(0.25), "▼", 14, False, PURPLE, PP_ALIGN.CENTER)

# Tech stack
txt(s, Inches(1.5), Inches(6.7), Inches(10.3), Inches(0.4),
    "Tech: Spring Boot 3.4 | ActiveMQ | JGraphT | Llama 3.2 + QLoRA | Ollama | AWS Bedrock | H2 | Python + Java | Jenkins | Docker",
    10, False, DIM, PP_ALIGN.CENTER)
slide_num(s, 3)

# ═══════════════════════════════════════════════════════════
# SLIDE 4: Our Own LLM
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(4), Inches(0.3), "CUSTOM LLM", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "KAVACH Compliance LLM - Our Own Model", 32, True, WHITE)

# Left: Fine-Tuning Pipeline
card(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(5.0), LIGHT_PURPLE)
multi(s, Inches(1.1), Inches(2.1), Inches(5.3), Inches(4.7), [
    ("Fine-Tuning Pipeline - How We Built It", LIGHT_PURPLE, True, 16),
    ("", GRAY, False, 8),
    ("Step 1: Collect training pairs", LIGHT_PURPLE, True, 13),
    ("Each example = \"here's code\" -> \"here's the compliance violation\"", GRAY, False, 12),
    ("Source: audit findings, Checkmarx reports, SOX evidence, MSA clauses", GRAY, False, 12),
    ("Volume: 500-1000 pairs (started with 15, expanding)", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("Step 2: Freeze model, train small adapters (QLoRA)", LIGHT_PURPLE, True, 13),
    ("Base model: Llama 3.2 (1.2B params) - frozen at 99.7%", GRAY, False, 12),
    ("Add 4M trainable adapter params on top", GRAY, False, 12),
    ("Train only adapters on our data (~2-6 hours, 1 GPU, ~$50-100)", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("Step 3: Merge, export, deploy", LIGHT_PURPLE, True, 13),
    ("Merge adapters -> export to Ollama format", GRAY, False, 12),
    ("Result: kavach-compliance-v1 - runs on our infra", GREEN, False, 12),
    ("No client code sent to external APIs", GREEN, False, 12),
], size=12, spacing=3)

# Right: Comparison table
card(s, Inches(6.9), Inches(1.9), Inches(5.8), Inches(5.0), GREEN)
multi(s, Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.7), [
    ("Why \"Our Own\" - Not Just a Wrapper", GREEN, True, 16),
    ("", GRAY, False, 8),
    ("Who sees client code?", DIM, True, 12),
    ("  GPT/Claude: OpenAI / Anthropic", RED, False, 12),
    ("  KAVACH: Nobody (runs on our GPU)", GREEN, False, 12),
    ("", GRAY, False, 4),
    ("Cost @ 5K events/day?", DIM, True, 12),
    ("  GPT/Claude: ~$6,000/month", RED, False, 12),
    ("  KAVACH: ~$1,000/month (1 GPU)", GREEN, False, 12),
    ("", GRAY, False, 4),
    ("Where does it learn?", DIM, True, 12),
    ("  GPT/Claude: General internet data", RED, False, 12),
    ("  KAVACH: OUR audit findings & controls", GREEN, False, 12),
    ("", GRAY, False, 4),
    ("Who owns the model?", DIM, True, 12),
    ("  GPT/Claude: OpenAI / Anthropic", RED, False, 12),
    ("  KAVACH: Hexaware (our IP)", GREEN, False, 12),
], size=12, spacing=2)
slide_num(s, 4)

# ═══════════════════════════════════════════════════════════
# SLIDE 5: Graph RAG
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(6), Inches(0.3), "KNOWLEDGE GRAPH + GraphRAG", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "Compliance Knowledge Graph with Graph RAG", 32, True, WHITE)

card(s, Inches(0.8), Inches(1.9), Inches(5.8), Inches(5.0), AMBER)
multi(s, Inches(1.1), Inches(2.1), Inches(5.3), Inches(4.7), [
    ("How GraphRAG Works in KAVACH", AMBER, True, 16),
    ("", GRAY, False, 6),
    ("1. Compliance Knowledge Graph", LIGHT_PURPLE, True, 13),
    ("All regulations, controls, systems modeled as connected nodes.", GRAY, False, 12),
    ("Edges: violates, requires, triggers, penalizes.", GRAY, False, 12),
    ("", GRAY, False, 4),
    ("2. Causal Path Traversal", LIGHT_PURPLE, True, 13),
    ("On event: traverse graph for all downstream impacts.", GRAY, False, 12),
    ("E.g.: SQL injection -> PCI-DSS 6.5 -> MSA SLA -> $50K penalty", GRAY, False, 12),
    ("", GRAY, False, 4),
    ("3. RAG (Retrieval-Augmented Generation)", LIGHT_PURPLE, True, 13),
    ("Retrieved paths passed as context to LLM.", GRAY, False, 12),
    ("Enables precise, regulation-aware reasoning.", GRAY, False, 12),
    ("", GRAY, False, 4),
    ("4. Multi-Domain Impact Output", LIGHT_PURPLE, True, 13),
    ("Structured response: domains, regulations, severity, actions.", GRAY, False, 12),
    ("All within 2 seconds of the triggering event.", GREEN, False, 12),
], size=12, spacing=3)

card(s, Inches(6.9), Inches(1.9), Inches(5.8), Inches(5.0), BORDER)
multi(s, Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.7), [
    ("Example: How One Finding Cascades", WHITE, True, 16),
    ("", GRAY, False, 6),
    ("A developer introduces SQL injection in loan search.", GRAY, False, 12),
    ("KAVACH traces the full impact through the graph:", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("SQL Injection Found", RED, True, 13),
    ("  |-- violates PCI-DSS 6.5 (secure coding)", GRAY, False, 12),
    ("  |-- violates OWASP A03 (injection prevention)", GRAY, False, 12),
    ("  |-- triggers MSA 7.2 SLA (48-hour clock starts)", AMBER, False, 12),
    ("  |    '-- penalty: $50,000/incident", RED, False, 12),
    ("  |-- requires Code Review (before fix ships)", GRAY, False, 12),
    ("  |    '-- satisfies ITGC-CM-06 (SOX)", GREEN, False, 12),
    ("  '-- blocks Deployment (gate: BLOCKED)", RED, True, 12),
    ("", GRAY, False, 6),
    ("What AI receives: traced paths + actual code +", DIM, False, 11),
    ("client's contract terms -> complete assessment.", DIM, False, 11),
], size=12, spacing=2)
slide_num(s, 5)

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Multi-Agent System
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "AGENTIC ARCHITECTURE", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "5 Autonomous Agents - Perceive, Reason, Act", 32, True, WHITE)

agents = [
    ("Digital Twin Agent", "Per-client compliance state", "Scores, risks, frameworks, SLAs", PURPLE),
    ("Chain Reactor Agent", "Cross-domain causal reasoning", "1 event -> 6 domains analyzed", AMBER),
    ("Audit Narrator Agent", "Auto-generates evidence", "LLM writes audit narratives", GREEN),
    ("Drift Sentinel Agent", "Silent degradation detection", "Access creep, config drift", RED),
    ("Obligation Parser Agent", "Contract -> machine rules", "MSA text -> enforceable policies", LIGHT_PURPLE),
    ("Control Ingestion Agent", "Add new frameworks via AI", "Zero code changes needed", AMBER),
]
for i, (name, line1, line2, color) in enumerate(agents):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.2)
    y = Inches(2.0 + row * 2.6)
    card(s, x, y, Inches(3.9), Inches(2.2), color)
    multi(s, x + Pt(14), y + Pt(10), Inches(3.5), Inches(2.0), [
        (name, color, True, 15),
        ("", GRAY, False, 4),
        (line1, WHITE, False, 13),
        (line2, GRAY, False, 12),
    ], spacing=4)
slide_num(s, 6)

# ═══════════════════════════════════════════════════════════
# SLIDE 7: Tech Stack
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "TECHNOLOGY STACK", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "Enterprise-Grade Tech Stack", 32, True, WHITE)

techs = [
    ("AI / LLM", "Llama 3.2 (fine-tuned)\nQLoRA Training\nOllama Runtime\nAWS Bedrock (Claude)"),
    ("Backend", "Spring Boot 3.4\nActiveMQ (JMS)\nH2 Database\nJPA / Hibernate"),
    ("Graph / RAG", "JGraphT\nCausal Path Analysis\nGraphRAG Retrieval\nKnowledge Graph"),
    ("Frontend", "Real-time Dashboard\nAuto-refresh (no reload)\n5 Interactive Views\nClick-to-expand"),
    ("Code Analysis", "AI Scanner (Ollama)\nEntropy-based secrets\nSAST integration\nFile System Watcher"),
    ("Event Pipeline", "Webhook ingestion\nJMS message queue\nEvent-driven agents\nAsync processing"),
    ("Microservices", "Gradle Multi-project\nAlpha/Beta/Gamma\nSelective Deploy\nJenkins Pipeline"),
    ("Cloud / Infra", "AWS (Bedrock, S3)\nDocker / K8s ready\nZero external deps\nSelf-contained demo"),
]
for i, (title, items) in enumerate(techs):
    col = i % 4
    row = i // 4
    x = Inches(0.8 + col * 3.15)
    y = Inches(2.0 + row * 2.7)
    card(s, x, y, Inches(2.9), Inches(2.4), PURPLE)
    multi(s, x + Pt(12), y + Pt(8), Inches(2.6), Inches(2.2), [
        (title, LIGHT_PURPLE, True, 14),
        ("", GRAY, False, 4),
    ] + [(line, GRAY, False, 11) for line in items.split('\n')], spacing=3)
slide_num(s, 7)

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Chain Reaction
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "KEY DIFFERENTIATOR", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "One Commit -> Six Compliance Domains", 32, True, WHITE)

# Left card: Scenario
card(s, Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.2), AMBER)
multi(s, Inches(1.1), Inches(2.0), Inches(5.5), Inches(5.0), [
    ("Scenario: Developer Pushes Code to Rate Calculation", WHITE, True, 14),
    ("", GRAY, False, 4),
    ("A routine commit modifies ARM interest rate logic -", GRAY, False, 12),
    ("touches financially-significant code and borrower PII.", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("What changed: Rate cap + eligibility logic", AMBER, False, 12),
    ("Impact: Determines borrower's rate", AMBER, False, 12),
    ("SAST: Race condition (HIGH)", RED, False, 12),
    ("Data: Income, credit score (PII)", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("Chain Reactor fires - 6 domains in <2 seconds:", GREEN, True, 13),
    ("", GRAY, False, 4),
    ("  [SOX] Dual approval required for financial system", WHITE, False, 12),
    ("  [Security] SAST finding blocks deployment", WHITE, False, 12),
    ("  [TILA/Reg Z] APR accuracy validation needed", WHITE, False, 12),
    ("  [Fair Lending] Disparate impact analysis required", WHITE, False, 12),
    ("  [Contractual] Client notification per MSA", WHITE, False, 12),
    ("  [PCI-DSS] Code review required before release", WHITE, False, 12),
], spacing=3)

# Right card: Competitor comparison
card(s, Inches(7.1), Inches(1.8), Inches(5.6), Inches(5.2), BORDER)
multi(s, Inches(7.4), Inches(2.0), Inches(5.1), Inches(5.0), [
    ("Why Existing Tools Fall Short", WHITE, True, 14),
    ("", GRAY, False, 6),
    ("Vanta / Drata ($25-50K/yr)", DIM, True, 12),
    ("  Single-company only. No per-client context.", GRAY, False, 11),
    ("", GRAY, False, 3),
    ("ServiceNow GRC ($100-300K/yr)", DIM, True, 12),
    ("  Aggregation, no causal reasoning across domains.", GRAY, False, 11),
    ("", GRAY, False, 3),
    ("Checkmarx / Snyk ($50-150K/yr)", DIM, True, 12),
    ("  Finds vulns, no SOX/TILA/contractual impact.", GRAY, False, 11),
    ("", GRAY, False, 3),
    ("Fieldguide ($30-80K/yr)", DIM, True, 12),
    ("  Helps auditors organize, doesn't generate evidence.", GRAY, False, 11),
    ("", GRAY, False, 6),
    ("KAVACH AI (~$1K/mo infra)", GREEN, True, 13),
    ("  Cross-domain causal reasoning + auto evidence.", GREEN, False, 12),
    ("  Covers ALL gaps above in one platform.", GREEN, False, 12),
    ("", GRAY, False, 4),
    ("Clients spend $200-500K/yr on 3-4 tools and STILL", DIM, False, 11),
    ("find gaps in audit. KAVACH replaces them all.", DIM, False, 11),
], spacing=2)
slide_num(s, 8)

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Dynamic Control Ingestion
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(6), Inches(0.3), "ADAPTABILITY - ZERO CODE CHANGES", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "New Regulation Drops? KAVACH Adapts Instantly.", 32, True, WHITE)
txt(s, Inches(0.8), Inches(1.5), Inches(12), Inches(0.4),
    "Traditional: weeks/months to implement new framework.  KAVACH: paste regulation text -> enforced in seconds.", 12, False, GRAY)

# Left: How it works
card(s, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.8), PURPLE)
multi(s, Inches(1.1), Inches(2.3), Inches(5.3), Inches(4.5), [
    ("How It Works - Real Example", LIGHT_PURPLE, True, 15),
    ("", GRAY, False, 4),
    ("Scenario: DORA regulation takes effect. Client EuroLend", GRAY, False, 12),
    ("is now subject to it. Enforce it TODAY.", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("1. Compliance officer pastes regulation text", WHITE, True, 13),
    ("   Copies Article 19 directly from DORA - no formatting", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("2. Our LLM reads and understands legal language", WHITE, True, 13),
    ("   Determines: CRITICAL, REGULATORY, 4h SLA, blocking", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("3. Policy added to live engine - immediately active", WHITE, True, 13),
    ("   No developer, no code change, no restart.", GRAY, False, 11),
    ("", GRAY, False, 6),
    ("Time: regulation published -> enforced = MINUTES", GREEN, True, 13),
    ("", GRAY, False, 4),
    ("Works for ANY industry:", DIM, False, 11),
    ("Mortgage (SOX, TILA) | Healthcare (HIPAA) | Insurance", DIM, False, 10),
    ("Retail (PCI-DSS) | Airlines (FAA) | Public (FedRAMP)", DIM, False, 10),
], spacing=2)

# Right: What AI produces
card(s, Inches(6.9), Inches(2.1), Inches(5.8), Inches(4.8), BORDER)
multi(s, Inches(7.2), Inches(2.3), Inches(5.3), Inches(4.5), [
    ("What Happens Under the Hood", WHITE, True, 15),
    ("", GRAY, False, 4),
    ("Compliance officer pastes DORA Article 19 (legal text).", GRAY, False, 12),
    ("LLM asks: What type? How severe? What triggers? Block?", GRAY, False, 12),
    ("", GRAY, False, 6),
    ("INPUT (Legal Text):", RED, True, 12),
    ('"Financial entities shall report major ICT incidents', GRAY, False, 11),
    (' to competent authority within 4 hours..."', GRAY, False, 11),
    ("", GRAY, False, 4),
    ("OUTPUT (Enforceable Policy):", GREEN, True, 12),
    ("  Domain: REGULATORY", GRAY, False, 11),
    ("  Severity: CRITICAL", GRAY, False, 11),
    ("  SLA: 4 hours", GRAY, False, 11),
    ("  Blocking: Yes", GRAY, False, 11),
    ("", GRAY, False, 6),
    ("Business impact for Hexaware:", LIGHT_PURPLE, True, 12),
    ("  Client adopts framework -> enforce same day", WHITE, False, 11),
    ("  No developer sprint -> saves 2-4 weeks", WHITE, False, 11),
    ("  Compliance team self-serves -> no eng dependency", WHITE, False, 11),
    ("  Scales to 100+ clients with different frameworks", WHITE, False, 11),
], spacing=2)
slide_num(s, 9)

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Business Value
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); gradient_bar(s)
txt(s, Inches(0.8), Inches(0.5), Inches(5), Inches(0.3), "BUSINESS IMPACT", 11, True, PURPLE)
txt(s, Inches(0.8), Inches(0.9), Inches(12), Inches(0.6), "Value for Hexaware & Clients", 32, True, WHITE)

# Left: Quantified Value
card(s, Inches(0.8), Inches(1.8), Inches(5.8), Inches(5.2), PURPLE)
multi(s, Inches(1.1), Inches(2.0), Inches(5.3), Inches(5.0), [
    ("Quantified Value - How We Calculated", LIGHT_PURPLE, True, 15),
    ("", GRAY, False, 6),
    ("70% Audit Prep Time Saved", LIGHT_PURPLE, True, 14),
    ("Auditors spend 3-4 weeks collecting evidence. KAVACH", GRAY, False, 11),
    ("generates it continuously. Auditors walk in -> done.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("$2-4M Annual Savings", GREEN, True, 14),
    ("Audit prep labor ($1.2M) + penalty avoidance ($800K)", GRAY, False, 11),
    ("+ reduced headcount (3-4 FTEs) + tool consolidation.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("$50K+ Penalty Avoided Per Incident", AMBER, True, 14),
    ("MSA 7.2: $50K per critical not fixed in 48h.", GRAY, False, 11),
    ("Avg client: 5-10 criticals/yr. 2-3 breach SLA = $100-150K.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("<2 Second Analysis Time", RED, True, 14),
    ("Current: compliance reviews 2-5 days after deploy.", GRAY, False, 11),
    ("KAVACH: at commit time. 6 domains in under 2 seconds.", GRAY, False, 11),
], spacing=2)

# Right: Why this wins
card(s, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.2), GREEN)
multi(s, Inches(7.2), Inches(2.0), Inches(5.3), Inches(5.0), [
    ("Why This Wins - Competitive Advantage", GREEN, True, 15),
    ("", GRAY, False, 6),
    ("Own LLM", LIGHT_PURPLE, True, 13),
    ("Client code never leaves our infra. Model improves", GRAY, False, 11),
    ("with every engagement - compounding advantage.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("Graph RAG + Multi-Agent", LIGHT_PURPLE, True, 13),
    ("5 agents + knowledge graph. One finding traces impact", GRAY, False, 11),
    ("across SOX + Security + TILA + ECOA + PCI + Contract.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("Domain-Agnostic / Per-Engagement", LIGHT_PURPLE, True, 13),
    ("Same platform: mortgage, healthcare, retail, airlines.", GRAY, False, 11),
    ("Each client gets own compliance profile.", GRAY, False, 11),
    ("", GRAY, False, 4),
    ("Predictive, Not Reactive", LIGHT_PURPLE, True, 13),
    ("Blocks violations at commit time. Auditors find zero", GRAY, False, 11),
    ("surprises because system already enforced every control.", GRAY, False, 11),
    ("", GRAY, False, 6),
    ("KAVACH AI", PURPLE, True, 16),
    ('"One commit. Six domains. Zero breaches."', DIM, False, 12),
], spacing=2)
slide_num(s, 10)

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
import os
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KAVACH_AI_Presentation.pptx")
prs.save(output)
print(f"Done! Saved: {output}")
print("10 slides with same content as HTML presentation.")
