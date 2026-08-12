#!/usr/bin/env python3
"""
Generate KAVACH AI presentation deck with professional styling.
Uses python-pptx library.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_BG = RGBColor(0x0A, 0x0E, 0x1A)
ACCENT_PURPLE = RGBColor(0x6C, 0x63, 0xFF)
ACCENT_LIGHT = RGBColor(0xB8, 0xB4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x88, 0x88, 0x99)
RED_ALERT = RGBColor(0xF4, 0x43, 0x36)
GREEN_OK = RGBColor(0x4C, 0xAF, 0x50)
ORANGE_WARN = RGBColor(0xFF, 0x98, 0x00)
DARK_CARD = RGBColor(0x16, 0x1B, 0x30)
CARD_BORDER = RGBColor(0x2A, 0x2F, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle shape as background element."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with specified styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=LIGHT_GRAY, spacing=1.2):
    """Add multiple lines of text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(font_size * spacing * 0.4)
    return txBox


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
set_slide_bg(slide, DARK_BG)

# Accent bar at top
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

# Product name
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "KAVACH AI", font_size=54, bold=True, color=WHITE)

# Subtitle
add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
             "Kontinuous Audit & Vulnerability Analysis for Compliant Hardening", font_size=28, color=ACCENT_LIGHT)

# Tagline
add_text_box(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.6),
             '"One commit. Six domains. Zero breaches."',
             font_size=20, color=MEDIUM_GRAY)

# Category badge
add_text_box(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
             "Agentic Arena 2026  |  Autonomous Security, Compliance & Audit Intelligence",
             font_size=14, color=MEDIUM_GRAY)

# Bottom accent bar
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_PURPLE)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "THE PROBLEM", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.6),
             "A single code commit can silently violate 6 compliance domains simultaneously",
             font_size=18, color=LIGHT_GRAY)

problems = [
    "SOX ITGC Controls — Unapproved change to a financial system",
    "Application Security — Unresolved vulnerability deployed to production",
    "TILA / RESPA Regulations — Incorrect loan rate calculation shipped",
    "Client Contractual SLAs — 48-hour remediation window missed ($50K penalty)",
    "Fair Lending (ECOA) — Pricing logic changed without disparate impact test",
    "GDPR / DORA — Data crossing regional boundaries without consent",
]

add_multi_text(slide, Inches(0.8), Inches(1.9), Inches(11), Inches(3.5),
               problems, font_size=17, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.3), Inches(11), Inches(1.0),
             "These domains live in separate tools, managed by separate teams, on separate audit cycles.\n"
             "Nobody connects the dots. Violations are discovered months later during audits.",
             font_size=16, color=ORANGE_WARN)

add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
             "Annual risk exposure: $5M - $10M per organization in regulated delivery",
             font_size=18, bold=True, color=RED_ALERT)


# ============================================================
# SLIDE 3: What is KAVACH AI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "WHAT IS KAVACH AI?", font_size=32, bold=True, color=WHITE)

desc_lines = [
    "A multi-agent AI platform that maintains a live compliance",
    "digital twin for each client engagement.",
    "",
    "When any event occurs (code commit, infrastructure change,",
    "new regulation, personnel change) — KAVACH propagates the",
    "compliance impact across ALL domains simultaneously:",
    "",
    "SOX  +  Security  +  Regulatory  +  Contractual  +  Fair Lending  +  Privacy  +  Audit",
]
add_multi_text(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(3.5),
               desc_lines, font_size=18, color=LIGHT_GRAY)

capabilities = [
    "Generates autonomous audit evidence in real-time",
    "Detects silent compliance drift before auditors find it",
    "Parses client contracts into machine-readable obligations",
    "Simulates compliance impact before deployment (predictive, not reactive)",
]
add_multi_text(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2.5),
               capabilities, font_size=16, color=GREEN_OK)


# ============================================================
# SLIDE 4: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "PLATFORM ARCHITECTURE", font_size=32, bold=True, color=WHITE)

# Orchestrator box
add_shape_bg(slide, Inches(2), Inches(1.3), Inches(9.3), Inches(0.9), DARK_CARD)
add_text_box(slide, Inches(2.2), Inches(1.4), Inches(9), Inches(0.8),
             "ORCHESTRATOR — Master Agent (Plans, Prioritizes, Delegates, Remembers)",
             font_size=15, bold=True, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

# Agent boxes
agents = [
    ("Digital\nTwin", Inches(1.5)),
    ("Chain\nReactor", Inches(3.8)),
    ("Audit\nNarrator", Inches(6.1)),
    ("Drift\nDetector", Inches(8.4)),
    ("Obligation\nParser", Inches(10.7)),
]
for name, left in agents:
    add_shape_bg(slide, left, Inches(2.6), Inches(1.9), Inches(1.2), DARK_CARD)
    add_text_box(slide, left, Inches(2.7), Inches(1.9), Inches(1.1),
                 name, font_size=13, bold=True, color=ACCENT_LIGHT, alignment=PP_ALIGN.CENTER)

# Integration layer
add_shape_bg(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.9), DARK_CARD)
add_text_box(slide, Inches(1.7), Inches(4.4), Inches(10), Inches(0.8),
             "INTEGRATION LAYER: Snyk | Checkmarx | Wiz | Qualys | ServiceNow | Jira | Git | AWS | Splunk",
             font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom description
add_text_box(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(1.5),
             "KAVACH does not replace any existing tool.\n"
             "It sits above them as the orchestration, reasoning, and action layer\n"
             "that makes them all work together intelligently.",
             font_size=16, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5: Unique Capability 1 - Chain Reaction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "CROSS-DOMAIN CHAIN REACTION", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "One code commit triggers impact analysis across all compliance domains simultaneously",
             font_size=16, color=MEDIUM_GRAY)

# Event trigger
add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.7), DARK_CARD)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(11), Inches(0.6),
             "EVENT: Developer modifies ARM rate calculation logic (commit a3f7b2c)",
             font_size=15, bold=True, color=ORANGE_WARN)

# Impact chain
chain_items = [
    ("SOX", "Change management documentation required (ITGC-CM-01, ITGC-CM-02)"),
    ("SECURITY", "SAST scan shows 1 High finding (PCI-DSS 6.5 violation)"),
    ("REGULATORY", "TILA Regulation Z — APR validation to 1/8% accuracy needed"),
    ("FAIR LENDING", "Pricing logic change requires disparate impact testing (ECOA)"),
    ("CONTRACTUAL", "Client notification required per MSA Section 7.4 (48h before deploy)"),
    ("AUDIT", "Full evidence trail must be generated for SOX-critical change"),
]

y_pos = 2.8
for domain, desc in chain_items:
    add_text_box(slide, Inches(1.2), Inches(y_pos), Inches(2.2), Inches(0.4),
                 domain, font_size=13, bold=True, color=ACCENT_PURPLE)
    add_text_box(slide, Inches(3.4), Inches(y_pos), Inches(9), Inches(0.4),
                 desc, font_size=13, color=LIGHT_GRAY)
    y_pos += 0.55

add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11), Inches(0.5),
             "No product on the market performs cross-domain causal compliance reasoning. They all operate in silos.",
             font_size=15, bold=True, color=ACCENT_LIGHT)


# ============================================================
# SLIDE 6: Unique Capability 2 - Pre-Deployment Simulation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "PRE-DEPLOYMENT COMPLIANCE SIMULATION", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Know the compliance impact BEFORE you ship — not months later during audit",
             font_size=16, color=MEDIUM_GRAY)

# Simulation output box
add_shape_bg(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(4.2), DARK_CARD)

sim_lines = [
    "COMPLIANCE SIMULATION RESULT",
    "",
    "If you deploy this change, Client A's compliance score drops from 94% to 87%",
    "",
    "Blocking Issues:",
    "   Unresolved SAST finding — SOX ITGC-SD-01 violation",
    "   Client notification not yet sent — MSA breach risk, $50K penalty",
    "   APR validation not executed — TILA regulatory risk",
    "",
    "Estimated financial risk if deployed now: $250,000",
    "Recommendation: Hold deployment. Resolution time: ~6 hours",
]
add_multi_text(slide, Inches(1.2), Inches(2.1), Inches(11), Inches(3.8),
               sim_lines, font_size=15, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
             "Every existing tool is reactive. KAVACH is the only platform that simulates consequences before action.",
             font_size=15, bold=True, color=ACCENT_LIGHT)


# ============================================================
# SLIDE 7: Unique Capability 3 - Audit Narrator
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "AUTONOMOUS AUDIT NARRATIVE GENERATION", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Real-time evidence generation from development activity — not manual reconstruction months later",
             font_size=16, color=MEDIUM_GRAY)

# Narrative example
add_shape_bg(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.8), DARK_CARD)

narrative_lines = [
    "AUTO-GENERATED AUDIT NARRATIVE (Sample)",
    "",
    "On June 28, 2026, PR #4521 was created by Developer A to modify the ARM",
    "rate calculation module (SOX-critical system).",
    "Business justification: CFPB Bulletin 2026-03 (Jira: MORT-1542).",
    "SAST scan by Checkmarx: 0 Critical, 1 High finding — resolved in 4 hours.",
    "Code reviewed by Developer B (segregation of duties: SATISFIED).",
    "Deployed to production June 29 with change ticket CHG-3302.",
    "Dual-approved by Release Manager C.",
    "Controls satisfied: ITGC-CM-01, ITGC-CM-02, ITGC-SD-01, PCI-DSS 6.5.1.",
]
add_multi_text(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(3.5),
               narrative_lines, font_size=14, color=LIGHT_GRAY)

add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11), Inches(0.8),
             "Audit prep time reduced by 70%. EY/Deloitte get pre-built evidence packages\n"
             "instead of spending weeks reconstructing timelines from Jira, Git, and ServiceNow.",
             font_size=15, bold=True, color=GREEN_OK)


# ============================================================
# SLIDE 8: Market Differentiation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "MARKET DIFFERENTIATION", font_size=32, bold=True, color=WHITE)

add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "Verified against 30+ products — this capability does not exist today",
             font_size=16, color=MEDIUM_GRAY)

# Competitor comparison
competitors = [
    ("Vanta / Drata / Scrut", "Evidence collection for one company", "No multi-client, no reasoning"),
    ("ServiceNow GRC", "Risk register, control tracking", "Static, no causal chain analysis"),
    ("Panther / CrowdStrike", "SOC automation, incident response", "Security-only, no SOX/regulatory"),
    ("Checkmarx / Snyk / Wiz", "Detect vulnerabilities", "No business impact, no audit evidence"),
    ("Fieldguide / AuditBoard", "Auditor workflow management", "Don't generate evidence autonomously"),
    ("Regology / Wolters Kluwer", "Track regulatory changes", "Don't simulate per-client blast radius"),
]

# Header
y = 1.6
add_text_box(slide, Inches(0.8), Inches(y), Inches(3.2), Inches(0.4),
             "Product", font_size=12, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, Inches(4.0), Inches(y), Inches(4.0), Inches(0.4),
             "What They Do", font_size=12, bold=True, color=ACCENT_PURPLE)
add_text_box(slide, Inches(8.2), Inches(y), Inches(4.5), Inches(0.4),
             "What They Miss", font_size=12, bold=True, color=ACCENT_PURPLE)

y = 2.1
for product, does, misses in competitors:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(3.2), Inches(0.45),
                 product, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, Inches(4.0), Inches(y), Inches(4.0), Inches(0.45),
                 does, font_size=13, color=LIGHT_GRAY)
    add_text_box(slide, Inches(8.2), Inches(y), Inches(4.5), Inches(0.45),
                 misses, font_size=13, color=RED_ALERT)
    y += 0.62

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "KAVACH sits ABOVE all these tools as the orchestration and reasoning layer.\n"
             "It does not replace them — it makes them work together intelligently.",
             font_size=15, bold=True, color=ACCENT_LIGHT)


# ============================================================
# SLIDE 9: Business Value
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "BUSINESS VALUE", font_size=32, bold=True, color=WHITE)

# Value metrics
metrics = [
    ("70%", "Audit preparation\ntime reduced"),
    ("$2M", "SLA breach penalties\navoided per year"),
    ("80%", "Material weakness\nfindings prevented"),
    ("175%", "Year 1 ROI\n(conservative)"),
]

x_pos = 0.8
for value, label in metrics:
    add_shape_bg(slide, Inches(x_pos), Inches(1.4), Inches(2.8), Inches(2.2), DARK_CARD)
    add_text_box(slide, Inches(x_pos), Inches(1.6), Inches(2.8), Inches(1.0),
                 value, font_size=40, bold=True, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(x_pos), Inches(2.7), Inches(2.8), Inches(0.8),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += 3.1

# Additional value points
value_lines = [
    "Developer compliance overhead reduced by 20-30%",
    "New revenue stream: premium managed service at $20K-$50K/month per client",
    "Competitive differentiation: no competitor in IT services has this capability",
    "Client retention: real-time compliance visibility creates trust and switching costs",
    "Audit relationship transformed: auditors become collaborators, not adversaries",
]
add_multi_text(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(3.0),
               value_lines, font_size=16, color=LIGHT_GRAY)


# ============================================================
# SLIDE 10: 5 Autonomous Agents
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "FIVE AUTONOMOUS AGENTS", font_size=32, bold=True, color=WHITE)

agents_data = [
    ("Digital Twin Agent", "Builds and maintains live compliance score per engagement from all tool feeds"),
    ("Chain Reactor Agent", "Propagates any event across all compliance domains with causal reasoning"),
    ("Audit Narrator Agent", "Generates continuous, audit-ready evidence narratives from dev activity"),
    ("Drift Detector Agent", "Finds silent compliance degradation: expired certs, orphan access, config drift"),
    ("Obligation Parser Agent", "Extracts machine-readable rules from contracts using LLM"),
]

y = 1.3
for agent_name, agent_desc in agents_data:
    add_shape_bg(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(1.0), DARK_CARD)
    add_text_box(slide, Inches(1.2), Inches(y + 0.1), Inches(3.5), Inches(0.5),
                 agent_name, font_size=15, bold=True, color=ACCENT_LIGHT)
    add_text_box(slide, Inches(1.2), Inches(y + 0.5), Inches(10.8), Inches(0.5),
                 agent_desc, font_size=14, color=LIGHT_GRAY)
    y += 1.15


# ============================================================
# SLIDE 11: Tech Stack & Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "TECH STACK & ROADMAP", font_size=32, bold=True, color=WHITE)

# Roadmap
phases = [
    ("Phase 1 (8 weeks)", "Core Twin + Chain Reactor, 3 tool integrations, 1 pilot client"),
    ("Phase 2 (12 weeks)", "Obligation Parser + Audit Narrator, Jira/GitHub/AWS integration"),
    ("Phase 3 (16 weeks)", "Multi-client command center, Drift Detection, client portal"),
    ("Phase 4 (20 weeks)", "Full multi-tenant platform, API marketplace, white-label ready"),
]

y = 1.2
for phase, desc in phases:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(3.0), Inches(0.4),
                 phase, font_size=14, bold=True, color=ACCENT_PURPLE)
    add_text_box(slide, Inches(3.8), Inches(y), Inches(9), Inches(0.4),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += 0.55

# Tech stack
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.5),
             "TECHNOLOGY", font_size=14, bold=True, color=ACCENT_LIGHT)

tech_items = [
    "Agent Framework: LangGraph (stateful multi-agent workflows)",
    "LLM: Claude / GPT-4o (contract parsing, narrative generation)",
    "Graph DB: Neo4j (obligation-control-tool relationships)",
    "Vector DB: Pinecone (regulatory text embeddings)",
    "Integrations: Snyk, Checkmarx, Wiz, ServiceNow, Jira, GitHub APIs",
    "Cloud: AWS multi-account (per-client isolation)",
    "Dashboard: React + D3.js (real-time visualization)",
    "Observability: LangSmith (agent decision tracing)",
]
add_multi_text(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(3.0),
               tech_items, font_size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 12: Why This Wins (Closing)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
             "WHY THIS WINS", font_size=32, bold=True, color=WHITE)

criteria = [
    ("Does it exist in the market?", "NO — verified against 30+ products"),
    ("Is it a real problem?", "YES — $5-10M annual risk in regulated delivery"),
    ("Is it agentic?", "YES — 5 autonomous agents with reasoning and action"),
    ("Is it feasible?", "YES — working prototype built, standard API integrations"),
    ("Does it generate revenue?", "YES — premium service + platform licensing"),
    ("Is it defensible?", "YES — multi-client data, contract corpus, domain expertise"),
]

y = 1.3
for question, answer in criteria:
    add_text_box(slide, Inches(0.8), Inches(y), Inches(5.0), Inches(0.5),
                 question, font_size=16, color=LIGHT_GRAY)
    color = GREEN_OK if "YES" in answer or "NO —" in answer else LIGHT_GRAY
    add_text_box(slide, Inches(6.0), Inches(y), Inches(6.5), Inches(0.5),
                 answer, font_size=16, bold=True, color=color)
    y += 0.6

# Closing quote
add_shape_bg(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.8), DARK_CARD)
add_text_box(slide, Inches(1.2), Inches(5.4), Inches(11), Inches(1.5),
             "Every IT services company has Checkmarx. Every one has ServiceNow.\n"
             "Every one has Wiz. None of them have a brain that connects all these tools,\n"
             "reasons across compliance domains, and tells you — before you deploy —\n"
             "that this one code commit will violate three different clients' contracts\n"
             "in three different ways.\n\n"
             "That brain is KAVACH AI.",
             font_size=15, color=LIGHT_GRAY)


# ============================================================
# SLIDE 13: Thank You / Contact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), ACCENT_PURPLE)
add_shape_bg(slide, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08), ACCENT_PURPLE)

add_text_box(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(1.2),
             "KAVACH AI", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
             "Kontinuous Audit & Vulnerability Analysis for Compliant Hardening", font_size=24, color=ACCENT_LIGHT,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.6),
             '"One commit. Six domains. Zero breaches."',
             font_size=20, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
             "Working prototype ready for demo",
             font_size=16, color=GREEN_OK, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(6.4), Inches(10), Inches(0.5),
             "Agentic Arena 2026",
             font_size=14, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KAVACH_AI_Presentation.pptx")
prs.save(output_path)
print(f"\nPresentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("\nDone! Open the .pptx file in PowerPoint or Google Slides.")
