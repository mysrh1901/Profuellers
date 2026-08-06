#!/usr/bin/env python3
"""Generate final KAVACH AI presentation for panel review."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Colors
BG = RGBColor(0x0A, 0x0E, 0x1A)
PURPLE = RGBColor(0x6C, 0x63, 0xFF)
LIGHT_PURPLE = RGBColor(0xB8, 0xB4, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
DIM = RGBColor(0x77, 0x77, 0x88)
RED = RGBColor(0xF4, 0x43, 0x36)
GREEN = RGBColor(0x4C, 0xAF, 0x50)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
CARD = RGBColor(0x16, 0x1B, 0x30)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

def bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = PURPLE; s.line.fill.background()

def box(slide, l, t, w, h, color=CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = "Segoe UI"
    p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines, size=15, color=GRAY):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(size); p.font.color.rgb = color; p.font.name = "Segoe UI"
        p.space_after = Pt(4)

# SLIDE 1: Title
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(1.5), Inches(2.0), Inches(10), Inches(1), "KAVACH AI", 54, True, WHITE, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(3.3), Inches(10), Inches(0.7), "Knowledge-driven Audit, Vulnerability Analysis & Compliance Health", 26, False, LIGHT_PURPLE, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(4.3), Inches(10), Inches(0.5), '"One commit. Six domains. Zero breaches."', 18, False, DIM, PP_ALIGN.CENTER)
txt(s, Inches(1.5), Inches(5.8), Inches(10), Inches(0.4), "Agentic Arena 2026 | Autonomous Security, Compliance & Audit Intelligence", 13, False, DIM, PP_ALIGN.CENTER)

# SLIDE 2: Problem Statement
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "PROBLEM STATEMENT", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5), "IT services companies face silent multi-domain compliance violations every day", 16, False, DIM)
multi(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), [
    "THE SCENARIO:",
    "A developer commits code changing loan rate calculation logic.",
    "",
    "WHAT ACTUALLY HAPPENS (today):",
    "  - Checkmarx flags a SAST finding → AppSec team creates a Jira ticket",
    "  - Code deploys to production anyway (nobody connects SAST to SOX)",
    "  - 3 months later, EY audit asks: Where is your SOX change management evidence?",
    "  - Where is the TILA rate validation? The fair lending impact test?",
    "  - Result: Material weakness. $50K-$500K penalty. Client trust damaged.",
    "",
    "WHY THIS HAPPENS:",
    "  - SOX team uses ServiceNow. Security team uses Checkmarx. Regulatory team uses spreadsheets.",
    "  - These tools never talk to each other.",
    "  - Nobody maps the causal chain: 1 code change → 6 compliance domains affected.",
    "",
    "ANNUAL RISK EXPOSURE: $5M - $10M per organization in regulated delivery"
], 14, GRAY)
txt(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4), "This is not hypothetical. This happens every quarter across regulated engagements.", 14, True, ORANGE)

# SLIDE 3: Solution
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "SOLUTION: KAVACH AI", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5), "An autonomous agent that sits on the ecosystem, monitors every change, and triggers all compliance checks", 15, False, DIM)
multi(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "WHAT IT IS:",
    "  A policy-driven multi-agent system that monitors Git, Jenkins, Jira, Docker,",
    "  AWS — and when any change happens, evaluates ALL compliance domains simultaneously.",
    "",
    "HOW IT WORKS:",
    "  1. Event occurs (code commit, Jenkins build, infra change, access change)",
    "  2. Event published to message queue (ActiveMQ/SQS)",
    "  3. Chain Reactor Agent picks it up, builds context (what changed + which client)",
    "  4. Policy Engine evaluates 15+ policies across all applicable domains",
    "  5. Results: which domains are affected, what actions are needed, what's blocking",
    "  6. Audit Narrator Agent auto-generates evidence narrative",
    "  7. Dashboard updates in real-time. Deployment blocked if needed.",
    "",
    "KEY DESIGN DECISIONS:",
    "  - Policy-driven (not hardcoded) — add new vertical by adding policies, no code change",
    "  - Hybrid AI: deterministic rules for auditability + Claude/Bedrock for intelligence",
    "  - Domain-agnostic: same agents work for Mortgage, Healthcare, Retail, Manufacturing",
    "  - Zero vendor lock-in: all open source (Spring Boot, ActiveMQ, H2)"
], 14, GRAY)

# SLIDE 4: 8 Areas of Compliance Check
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "8 COMPLIANCE DOMAINS CHECKED PER EVENT", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "All triggered simultaneously from a single event — no silo, no gap", 14, False, DIM)

domains = [
    ("1. SOX ITGC", "Dual approval, segregation of duties, change documentation, test evidence"),
    ("2. Application Security", "SAST/DAST results, open CVEs, deployment gate (blocking)"),
    ("3. Regulatory", "TILA rate accuracy, RESPA disclosures, GDPR residency, DORA reporting"),
    ("4. Fair Lending", "Disparate impact when pricing or eligibility logic changes (ECOA)"),
    ("5. Contractual (MSA/SLA)", "Client notification deadlines, SLA timers, penalty triggers"),
    ("6. Privacy", "PII processing detection, consent gaps, cross-border transfers"),
    ("7. Infrastructure", "Cloud misconfigs, IAM drift, encryption, backup, data residency"),
    ("8. Audit Evidence", "Auto-generates full narrative for every change (EY/Deloitte ready)"),
]

y = 1.4
for title, desc in domains:
    box(s, Inches(0.8), Inches(y), Inches(11.5), Inches(0.65), CARD)
    txt(s, Inches(1.0), Inches(y + 0.05), Inches(3.2), Inches(0.5), title, 13, True, PURPLE)
    txt(s, Inches(4.0), Inches(y + 0.05), Inches(8), Inches(0.5), desc, 12, False, GRAY)
    y += 0.72

txt(s, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3), "All 8 domains evaluated in < 1 second. Policies configurable per client engagement.", 12, False, GREEN)

# SLIDE 5: Domain Agnostic (Multi-Vertical)
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "DOMAIN-AGNOSTIC: WORKS IN ANY VERTICAL", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "Same agent, different policies. Adding a new industry = adding policy config, zero code changes.", 14, False, DIM)

verticals = [
    ("Mortgage / Banking", "SOX, TILA, RESPA, ECOA, HMDA, PCI-DSS, GLBA", "Rate accuracy, fair lending, change mgmt"),
    ("Healthcare / Pharma", "HIPAA, HITECH, FDA 21 CFR Part 11, SOX", "PHI protection, clinical system validation"),
    ("Insurance", "SOX, NAIC Model Laws, CCPA, State Regulations", "Claims logic, policyholder data, reserving"),
    ("Retail / E-commerce", "PCI-DSS, CCPA, GDPR, SOX", "Payment security, customer data, financial reporting"),
    ("Manufacturing", "ISO 27001, NIST, SOX, ITAR, Export Controls", "IP protection, supply chain, access control"),
    ("Public Sector", "FedRAMP, FISMA, CMMC, ITAR", "Authorization boundaries, data classification"),
]

y = 1.5
for vertical, frameworks, focus in verticals:
    box(s, Inches(0.8), Inches(y), Inches(11.5), Inches(0.78), CARD)
    txt(s, Inches(1.0), Inches(y + 0.03), Inches(3.0), Inches(0.35), vertical, 13, True, WHITE)
    txt(s, Inches(1.0), Inches(y + 0.38), Inches(5.5), Inches(0.35), frameworks, 10, False, PURPLE)
    txt(s, Inches(6.5), Inches(y + 0.15), Inches(5.5), Inches(0.5), focus, 12, False, GRAY)
    y += 0.88

txt(s, Inches(0.8), Inches(6.9), Inches(11), Inches(0.3), "Currently: 15 policies loaded across 9 verticals and 11 compliance domains.", 13, True, GREEN)

# SLIDE 6: Live Demo Flow
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "LIVE DEMO: REAL-TIME AGENT REACTION", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "Multiple event sources feed the same agent pipeline — all produce compliance results in real-time", 14, False, DIM)

multi(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5), [
    "EVENT SOURCES DEMONSTRATED:",
    "",
    "  [Git Hook]     Developer commits code → hook detects financial/PII changes → fires event",
    "  [Postman]      POST /api/events with custom JSON payload → agent processes immediately",
    "  [Jenkins]      POST /api/webhooks/simulate/jenkins-deploy-blocked → SAST gate triggered",
    "  [Docker]       POST /api/webhooks/simulate/docker-push → container scan policy triggered",
    "  [AWS Config]   POST /api/webhooks/simulate/aws-data-residency-violation → GDPR blocked",
    "  [Jira]         POST /api/webhooks/simulate/jira-ticket → change traceability logged",
    "",
    "DEMO FLOW:",
    "  1. Open dashboard at http://localhost:9090 — show clean state",
    "  2. Fire events via Postman or Git commit",
    "  3. Refresh dashboard — chain reaction results appear instantly",
    "  4. Show auto-generated audit narrative (no human wrote it)",
    "  5. Show compliance score degraded, blocking flags raised",
    "",
    "WHAT THE PANEL SEES:",
    "  - One event → multiple compliance domains triggered simultaneously",
    "  - Agent decided which policies apply based on client profile",
    "  - Evidence generated automatically — ready for auditor review",
    "  - Deployment blocked because SAST finding exists (SOX ITGC-SD-01)"
], 13, GRAY)

# SLIDE 7: Technical Architecture
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "TECHNICAL ARCHITECTURE", 30, True, WHITE)

# Draw architecture boxes
# Event Sources
box(s, Inches(0.5), Inches(1.2), Inches(2.5), Inches(3.5), CARD)
txt(s, Inches(0.6), Inches(1.3), Inches(2.3), Inches(0.4), "EVENT SOURCES", 10, True, PURPLE)
multi(s, Inches(0.6), Inches(1.7), Inches(2.3), Inches(3.0), [
    "Git Hooks", "Jenkins Webhooks", "Jira Webhooks", "Docker Registry", "AWS Config", "Postman / API"
], 11, GRAY)

# Message Queue
box(s, Inches(3.3), Inches(2.2), Inches(2.0), Inches(1.5), CARD)
txt(s, Inches(3.4), Inches(2.3), Inches(1.8), Inches(0.3), "MESSAGE QUEUE", 9, True, ORANGE)
txt(s, Inches(3.4), Inches(2.7), Inches(1.8), Inches(0.8), "ActiveMQ\n(Embedded)\n\nJMS Queue:\ncompliance-events", 10, False, GRAY)

# Agents
box(s, Inches(5.6), Inches(1.0), Inches(3.8), Inches(4.2), CARD)
txt(s, Inches(5.7), Inches(1.1), Inches(3.6), Inches(0.35), "AUTONOMOUS AGENTS", 10, True, GREEN)
multi(s, Inches(5.7), Inches(1.5), Inches(3.6), Inches(3.5), [
    "Chain Reactor Agent",
    "  - JMS Listener",
    "  - Builds event context",
    "  - Evaluates Policy Engine",
    "  - Determines blocking",
    "",
    "Audit Narrator Agent",
    "  - Generates evidence",
    "  - Calls Bedrock/Claude",
    "",
    "Policy Engine",
    "  - 15 policies, 9 verticals",
    "  - Domain-agnostic rules"
], 10, GRAY)

# Output
box(s, Inches(9.7), Inches(1.2), Inches(3.2), Inches(3.5), CARD)
txt(s, Inches(9.8), Inches(1.3), Inches(3.0), Inches(0.3), "OUTPUTS", 10, True, PURPLE)
multi(s, Inches(9.8), Inches(1.7), Inches(3.0), Inches(3.0), [
    "Chain Reaction Report",
    "  (domains affected)",
    "",
    "Audit Narrative",
    "  (auto-generated)",
    "",
    "Compliance Score",
    "  (updated real-time)",
    "",
    "Deployment Gate",
    "  (block/allow)",
    "",
    "Dashboard + APIs"
], 10, GRAY)

# Data layer
box(s, Inches(3.3), Inches(5.5), Inches(9.6), Inches(0.8), CARD)
txt(s, Inches(3.5), Inches(5.6), Inches(9.2), Inches(0.6), "DATA LAYER: H2 Database (demo) / PostgreSQL+RDS (production) | Client Profiles | Policies | Events | Results | Narratives", 10, False, DIM)

# Bedrock
box(s, Inches(0.5), Inches(5.5), Inches(2.5), Inches(0.8), CARD)
txt(s, Inches(0.6), Inches(5.6), Inches(2.3), Inches(0.6), "AWS Bedrock\n(Claude - Anthropic)", 10, True, LIGHT_PURPLE)

# Tech labels
txt(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.5), "Stack: Java 17 | Spring Boot 3.4 | ActiveMQ (embedded) | H2 (embedded) | AWS Bedrock SDK | Gradle | Zero external infra needed for demo", 11, False, DIM)

# SLIDE 8: Agent Intelligence Detail
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "AGENT INTELLIGENCE: HOW IT REASONS", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "Hybrid approach: deterministic rules for auditability + LLM for intelligence", 14, False, DIM)

multi(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5), [
    "RULE-BASED ENGINE (Deterministic):",
    "  Used for: compliance decisions",
    "",
    "  Why: Auditors need explainable logic.",
    "  Every decision traces to a specific",
    "  regulation or contract clause.",
    "",
    "  Examples:",
    "  - SOX: financial change → dual approval",
    "  - PCI: SAST finding → block deploy",
    "  - TILA: rate logic → APR validation",
    "  - MSA: financial change → notify client",
    "",
    "  Panel can inspect every rule in:",
    "  CompliancePolicyEngine.java",
], 13, GRAY)

multi(s, Inches(6.8), Inches(1.5), Inches(5.5), Inches(5.5), [
    "LLM / BEDROCK (Intelligent):",
    "  Used for: understanding + generation",
    "",
    "  Why: Some tasks need reasoning",
    "  beyond pattern matching.",
    "",
    "  Examples:",
    "  - Read code diff, understand intent",
    "  - Parse contract PDF into rules",
    "  - Generate natural language narrative",
    "  - Answer auditor questions",
    "",
    "  Model: Claude via Amazon Bedrock",
    "  (Hexaware = Anthropic authorized reseller)",
    "",
    "  Falls back to local logic when offline.",
], 13, GRAY)

txt(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4), "Key: The agent is agentic because it perceives → reasons → decides → acts autonomously. Not because it calls an LLM.", 13, True, LIGHT_PURPLE)

# SLIDE 9: Market Proof — This Does Not Exist
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "MARKET PROOF: THIS DOES NOT EXIST", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "Verified against 30+ products. No tool does cross-domain compliance reasoning from delivery events.", 14, False, DIM)

competitors = [
    ("Vanta / Drata / Scrut", "Evidence collection for one company", "No multi-client, no event reaction, no cross-domain"),
    ("ServiceNow GRC / SecOps", "Risk register, vulnerability tickets", "No causal reasoning, no policy engine, static"),
    ("Panther (Databricks)", "Agentic SOC, incident response", "Security only, no SOX/regulatory/contractual"),
    ("Checkmarx / Snyk / Wiz", "Detect vulnerabilities", "Detection only, no business impact or audit evidence"),
    ("Fieldguide / AuditBoard", "Auditor workflow tools", "For auditors, don't generate evidence from dev activity"),
    ("Regology / Wolters Kluwer", "Track regulatory changes", "Don't simulate impact per client engagement"),
    ("CrowdStrike / Splunk SOAR", "Security orchestration", "Incident response, not compliance reasoning"),
    ("GPT / Claude / Copilot (raw)", "Answer questions when asked", "Don't sit on ecosystem, don't react autonomously"),
]

y = 1.4
txt(s, Inches(0.8), Inches(y), Inches(3.5), Inches(0.3), "Product", 11, True, PURPLE)
txt(s, Inches(4.3), Inches(y), Inches(3.5), Inches(0.3), "What It Does", 11, True, PURPLE)
txt(s, Inches(7.8), Inches(y), Inches(5), Inches(0.3), "What It Misses (KAVACH fills this)", 11, True, PURPLE)
y += 0.4

for prod, does, misses in competitors:
    txt(s, Inches(0.8), Inches(y), Inches(3.5), Inches(0.35), prod, 10, True, WHITE)
    txt(s, Inches(4.3), Inches(y), Inches(3.5), Inches(0.35), does, 10, False, GRAY)
    txt(s, Inches(7.8), Inches(y), Inches(5), Inches(0.35), misses, 10, False, RED)
    y += 0.45

txt(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.6), "Conclusion: No product combines (a) ecosystem monitoring + (b) cross-domain policy evaluation +\n(c) autonomous evidence generation + (d) per-client context + (e) deployment blocking.", 12, True, LIGHT_PURPLE)

# SLIDE 10: Business Value
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "BUSINESS VALUE", 30, True, WHITE)

metrics = [("70%", "Audit prep\ntime reduced"), ("$2M+", "Penalties\navoided/year"), ("80%", "Findings\nprevented"), ("175%", "Year 1\nROI"), ("<1 sec", "Event to\nresult")]
x = 0.5
for val, label in metrics:
    box(s, Inches(x), Inches(1.1), Inches(2.3), Inches(1.8), CARD)
    txt(s, Inches(x), Inches(1.2), Inches(2.3), Inches(0.9), val, 32, True, PURPLE, PP_ALIGN.CENTER)
    txt(s, Inches(x), Inches(2.1), Inches(2.3), Inches(0.7), label, 12, False, GRAY, PP_ALIGN.CENTER)
    x += 2.5

multi(s, Inches(0.8), Inches(3.3), Inches(11.5), Inches(4.0), [
    "DIRECT SAVINGS:",
    "  - Audit preparation: from weeks of manual evidence gathering to always-ready (saves $2-4M/year)",
    "  - SLA breach prevention: agent tracks 48-hour / 7-day windows, alerts before deadline",
    "  - Penalty avoidance: $50K per critical vuln SLA miss × dozens of occurrences annually",
    "",
    "REVENUE OPPORTUNITY:",
    "  - Offer as premium managed service to clients: $20K-$50K/month per engagement",
    "  - Competitive differentiator: no TCS, Infosys, Wipro, Cognizant has this capability",
    "  - Platform licensing to other IT services companies (SaaS revenue)",
    "",
    "STRATEGIC VALUE:",
    "  - Transform auditor relationship: from adversarial to collaborative (evidence is always ready)",
    "  - Client retention: real-time compliance visibility creates trust and switching costs",
    "  - Category creation: 'Delivery Compliance Intelligence' doesn't exist yet — first mover wins",
    "  - Builds on Agentverse platform + Anthropic/Bedrock partnership"
], 13, GRAY)

# SLIDE 11: Code Overview for Panel
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
txt(s, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), "CODE: KEY FILES FOR PANEL REVIEW", 30, True, WHITE)
txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), "Spring Boot + ActiveMQ + H2 + AWS Bedrock — zero external infra, runs on any laptop", 14, False, DIM)

files = [
    ("CompliancePolicyEngine.java", "The brain — 15 configurable policies across 9 verticals, 11 domains"),
    ("ChainReactorAgent.java", "Autonomous agent — JMS listener, builds context, evaluates policies, acts"),
    ("AuditNarratorAgent.java", "Evidence generator — calls Bedrock/Claude or local logic"),
    ("BedrockClaudeService.java", "LLM integration — code analysis, narrative generation, contract parsing"),
    ("WebhookController.java", "Multi-source intake — Jenkins, Jira, Docker, AWS Config webhooks"),
    ("EventController.java", "REST API — accepts any event via POST, publishes to ActiveMQ"),
    ("HomeController.java", "Dashboard — real-time HTML view of twins, chain reactions, narratives"),
    ("post-commit (Git hook)", "Trigger — auto-fires event on every code commit, classifies changes"),
]

y = 1.4
for fname, desc in files:
    box(s, Inches(0.8), Inches(y), Inches(11.5), Inches(0.58), CARD)
    txt(s, Inches(1.0), Inches(y + 0.05), Inches(4.5), Inches(0.45), fname, 12, True, LIGHT_PURPLE)
    txt(s, Inches(5.5), Inches(y + 0.05), Inches(6.5), Inches(0.45), desc, 12, False, GRAY)
    y += 0.65

txt(s, Inches(0.8), Inches(6.8), Inches(11), Inches(0.4), "Total: ~800 lines of production-quality Java. All open source. No proprietary dependencies.", 13, True, GREEN)

# SLIDE 12: Closing
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s); bar(s)
box(s, Inches(1), Inches(1.5), Inches(11.3), Inches(4.5), CARD)
txt(s, Inches(1.5), Inches(1.8), Inches(10), Inches(0.6), "SUMMARY", 28, True, WHITE, PP_ALIGN.CENTER)

multi(s, Inches(1.5), Inches(2.5), Inches(10), Inches(3.5), [
    "What: An autonomous agent that monitors the delivery ecosystem and triggers",
    "      cross-domain compliance checks on every change — across 8 domains simultaneously.",
    "",
    "Why unique: No product in the market does this. Verified against 30+ tools.",
    "      Not Vanta. Not ServiceNow. Not Panther. Not Checkmarx. Not any raw LLM.",
    "",
    "Why it matters: Prevents $5-10M annual risk exposure. Reduces audit prep by 70%.",
    "      Creates new revenue as a managed service. Works in any vertical.",
    "",
    "What's built: Working prototype. Spring Boot + ActiveMQ + Bedrock.",
    "      Live demo available. Git commit → agents react → evidence generated in <1 second.",
], 14, GRAY, )

txt(s, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5), "KAVACH AI — \"One commit. Six domains. Zero breaches.\"", 18, True, PURPLE, PP_ALIGN.CENTER)


# SAVE
import os
path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KAVACH_AI_Final_Presentation.pptx")
prs.save(path)
print(f"Saved: {path}")
print(f"Slides: {len(prs.slides)}")
