#!/usr/bin/env python3
"""Clean, minimal PPT - white background, no clutter, flow diagrams."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
BLUE = RGBColor(0x00, 0x52, 0xCC)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xFE)
GRAY = RGBColor(0x5F, 0x63, 0x68)
LIGHT_GRAY = RGBColor(0xF1, 0xF3, 0xF4)
RED = RGBColor(0xD9, 0x3B, 0x3B)
GREEN = RGBColor(0x1E, 0x88, 0x55)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

def box(slide, l, t, w, h, color=LIGHT_GRAY, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align

def multi(slide, l, t, w, h, lines, size=14, color=GRAY):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)

def arrow(slide, l, t, w, h):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = BLUE
    s.line.fill.background()


# ===================== SLIDE 1: Title =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(2), Inches(2.5), Inches(9), Inches(1), "KAVACH AI", 52, True, BLACK, PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(3.7), Inches(9), Inches(0.6), "Kontinuous Audit & Vulnerability Analysis for Compliant Hardening", 24, False, GRAY, PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(4.8), Inches(9), Inches(0.5), "One commit. Six domains. Zero breaches.", 18, False, BLUE, PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(6.2), Inches(9), Inches(0.4), "Agentic Arena 2026", 14, False, GRAY, PP_ALIGN.CENTER)

# ===================== SLIDE 2: Problem =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "The Problem", 32, True, BLACK)

multi(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5), [
    "When a developer commits code in a regulated environment, that single change",
    "can violate multiple compliance domains at the same time:",
    "",
    "   SOX — unapproved change to financial system",
    "   Application Security — unresolved vulnerability goes to production",
    "   Regulatory (TILA/RESPA) — incorrect rate calculation",
    "   Contractual SLA — 48-hour remediation window missed",
    "   Fair Lending — pricing logic changed without impact testing",
    "   Privacy (GDPR/CCPA) — PII crosses regional boundary",
    "",
    "Today these are checked by different tools, owned by different teams,",
    "on different schedules. Nobody connects the chain reaction between them.",
    "",
    "Violations are discovered months later during audits.",
    "Cost: $50K - $500K per incident. Annual exposure: $5M - $10M.",
], 15, GRAY)

# ===================== SLIDE 3: Solution =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "The Solution", 32, True, BLACK)

multi(s, Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.0), [
    "An autonomous agent that sits on the delivery ecosystem,",
    "monitors every change, and triggers all compliance checks simultaneously.",
    "",
    "It does not replace any existing tool. It connects them into one intelligent layer.",
], 16, GRAY)

# Flow diagram
box(s, Inches(0.5), Inches(3.5), Inches(2.2), Inches(1.2), LIGHT_BLUE, BLUE)
txt(s, Inches(0.5), Inches(3.7), Inches(2.2), Inches(1.0), "Event Occurs\n\nCode commit\nInfra change\nAccess change", 11, False, BLACK, PP_ALIGN.CENTER)

arrow(s, Inches(2.8), Inches(3.9), Inches(0.6), Inches(0.3))

box(s, Inches(3.5), Inches(3.5), Inches(2.2), Inches(1.2), LIGHT_BLUE, BLUE)
txt(s, Inches(3.5), Inches(3.7), Inches(2.2), Inches(1.0), "Message Queue\n\nActiveMQ / SQS", 11, False, BLACK, PP_ALIGN.CENTER)

arrow(s, Inches(5.8), Inches(3.9), Inches(0.6), Inches(0.3))

box(s, Inches(6.5), Inches(3.5), Inches(2.5), Inches(1.2), LIGHT_BLUE, BLUE)
txt(s, Inches(6.5), Inches(3.7), Inches(2.5), Inches(1.0), "Chain Reactor\nAgent\n\nEvaluates all\npolicies", 11, False, BLACK, PP_ALIGN.CENTER)

arrow(s, Inches(9.1), Inches(3.9), Inches(0.6), Inches(0.3))

box(s, Inches(9.8), Inches(3.5), Inches(2.8), Inches(1.2), LIGHT_BLUE, BLUE)
txt(s, Inches(9.8), Inches(3.7), Inches(2.8), Inches(1.0), "Output\n\nBlock / Allow\nAudit Evidence\nScore Update", 11, False, BLACK, PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(5.2), Inches(11), Inches(0.4), "All 8 compliance domains evaluated in under 1 second per event.", 14, True, BLUE)


# ===================== SLIDE 4: 8 Domains =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "8 Compliance Domains Checked Per Event", 32, True, BLACK)
txt(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4), "All triggered simultaneously from a single change", 14, False, GRAY)

domains = [
    ("SOX ITGC", "Dual approval, segregation of duties, change documentation"),
    ("Application Security", "SAST/DAST findings, CVE tracking, deployment blocking"),
    ("Regulatory", "TILA rate accuracy, RESPA disclosures, GDPR data residency, DORA"),
    ("Fair Lending", "Disparate impact analysis when pricing or eligibility logic changes"),
    ("Contractual (MSA/SLA)", "Client notification deadlines, SLA timers, penalty triggers"),
    ("Privacy", "PII detection, consent gaps, cross-border data transfer"),
    ("Infrastructure", "Cloud misconfigs, IAM drift, encryption, backup retention"),
    ("Audit Evidence", "Auto-generated narrative for every change, ready on demand"),
]
y = 1.5
for i, (title, desc) in enumerate(domains, 1):
    box(s, Inches(0.8), Inches(y), Inches(11.5), Inches(0.63), LIGHT_GRAY)
    txt(s, Inches(1.0), Inches(y + 0.08), Inches(0.4), Inches(0.45), str(i), 12, True, BLUE, PP_ALIGN.CENTER)
    txt(s, Inches(1.4), Inches(y + 0.08), Inches(3.0), Inches(0.45), title, 13, True, BLACK)
    txt(s, Inches(4.5), Inches(y + 0.08), Inches(7.5), Inches(0.45), desc, 12, False, GRAY)
    y += 0.7

# ===================== SLIDE 5: Multi-Vertical =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Works Across Any Vertical", 32, True, BLACK)
txt(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4), "Same agents, different policies. Zero code changes to support a new industry.", 14, False, GRAY)

verticals = [
    ("Mortgage / Banking", "SOX, TILA, RESPA, ECOA, PCI-DSS, GLBA"),
    ("Healthcare / Pharma", "HIPAA, HITECH, FDA 21 CFR Part 11"),
    ("Insurance", "SOX, NAIC Model Laws, CCPA, State Regulations"),
    ("Retail / E-commerce", "PCI-DSS, CCPA, GDPR, SOX"),
    ("Manufacturing", "ISO 27001, NIST, ITAR, Export Controls"),
    ("Public Sector", "FedRAMP, FISMA, CMMC"),
]
y = 1.6
for vert, frameworks in verticals:
    box(s, Inches(0.8), Inches(y), Inches(11.5), Inches(0.7), LIGHT_GRAY)
    txt(s, Inches(1.2), Inches(y + 0.1), Inches(3.5), Inches(0.5), vert, 14, True, BLACK)
    txt(s, Inches(5.0), Inches(y + 0.1), Inches(7.0), Inches(0.5), frameworks, 13, False, BLUE)
    y += 0.8

txt(s, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4), "15 policies loaded across 9 verticals. Adding a new vertical = adding policy config only.", 13, True, GREEN)


# ===================== SLIDE 6: Architecture Flow =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Architecture", 32, True, BLACK)

# Row 1: Event Sources
txt(s, Inches(0.8), Inches(1.2), Inches(3), Inches(0.4), "Event Sources", 12, True, GRAY)
sources = ["Git Hooks", "Jenkins", "Jira", "Docker", "AWS Config", "Postman"]
y = 1.6
for src in sources:
    box(s, Inches(0.8), Inches(y), Inches(2.5), Inches(0.4), LIGHT_GRAY)
    txt(s, Inches(0.9), Inches(y + 0.02), Inches(2.3), Inches(0.35), src, 11, False, BLACK, PP_ALIGN.CENTER)
    y += 0.45

# Arrow to MQ
arrow(s, Inches(3.5), Inches(2.8), Inches(0.8), Inches(0.3))

# Message Queue
box(s, Inches(4.5), Inches(2.0), Inches(2.0), Inches(1.8), LIGHT_BLUE, BLUE)
txt(s, Inches(4.5), Inches(2.1), Inches(2.0), Inches(0.4), "Message Queue", 10, True, BLUE, PP_ALIGN.CENTER)
txt(s, Inches(4.5), Inches(2.5), Inches(2.0), Inches(1.2), "\nActiveMQ\n(embedded)\n\nJMS Queue", 10, False, GRAY, PP_ALIGN.CENTER)

# Arrow to Agents
arrow(s, Inches(6.7), Inches(2.8), Inches(0.8), Inches(0.3))

# Agents
box(s, Inches(7.7), Inches(1.3), Inches(2.8), Inches(3.5), LIGHT_BLUE, BLUE)
txt(s, Inches(7.7), Inches(1.4), Inches(2.8), Inches(0.4), "Agent Layer", 10, True, BLUE, PP_ALIGN.CENTER)
multi(s, Inches(7.9), Inches(1.8), Inches(2.5), Inches(3.0), [
    "Chain Reactor",
    "  Evaluates policies",
    "  Determines blocking",
    "",
    "Audit Narrator",
    "  Generates evidence",
    "  Calls LLM (Bedrock)",
    "",
    "Policy Engine",
    "  15 rules, 9 verticals",
], 10, GRAY)

# Arrow to Output
arrow(s, Inches(10.7), Inches(2.8), Inches(0.8), Inches(0.3))

# Output
box(s, Inches(11.7), Inches(1.8), Inches(1.4), Inches(2.5), LIGHT_BLUE, BLUE)
txt(s, Inches(11.7), Inches(1.9), Inches(1.4), Inches(0.3), "Output", 10, True, BLUE, PP_ALIGN.CENTER)
multi(s, Inches(11.8), Inches(2.2), Inches(1.3), Inches(2.0), [
    "Block/Allow",
    "deploy",
    "",
    "Audit",
    "narrative",
    "",
    "Score",
    "update",
    "",
    "Dashboard",
], 9, GRAY)

# Data layer at bottom
box(s, Inches(4.5), Inches(5.0), Inches(8.6), Inches(0.6), LIGHT_GRAY)
txt(s, Inches(4.6), Inches(5.1), Inches(8.4), Inches(0.4), "Data Layer: H2 (demo) / PostgreSQL (prod) — Events, Chain Reactions, Narratives, Client Profiles, Policies", 10, False, GRAY, PP_ALIGN.CENTER)

# Bedrock
box(s, Inches(7.7), Inches(5.8), Inches(2.8), Inches(0.6), LIGHT_GRAY)
txt(s, Inches(7.8), Inches(5.9), Inches(2.6), Inches(0.4), "AWS Bedrock (Claude) — optional LLM layer", 9, False, GRAY, PP_ALIGN.CENTER)


# ===================== SLIDE 7: Agent Flow (detailed) =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Agent Processing Flow", 32, True, BLACK)
txt(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4), "What happens when a code commit triggers the pipeline", 14, False, GRAY)

steps = [
    ("1", "Developer commits code", "Git hook detects file types\n(financial? PII? secrets?)"),
    ("2", "Event published to queue", "REST API persists event\nPublishes to ActiveMQ"),
    ("3", "Chain Reactor picks it up", "Builds context from event\n+ client profile"),
    ("4", "Policy Engine evaluates", "Checks all 15 policies\nagainst event context"),
    ("5", "Results determined", "Which domains hit?\nBlocking or not?"),
    ("6", "Audit Narrator generates", "Writes evidence narrative\n(LLM or local)"),
    ("7", "Score updated", "Compliance score recalculated\nDebt recomputed"),
    ("8", "Dashboard reflects", "Real-time update\nDrill-down available"),
]

x = 0.3
for num, title, desc in steps:
    box(s, Inches(x), Inches(1.7), Inches(1.5), Inches(2.5), LIGHT_BLUE, BLUE)
    txt(s, Inches(x), Inches(1.8), Inches(1.5), Inches(0.4), num, 16, True, BLUE, PP_ALIGN.CENTER)
    txt(s, Inches(x), Inches(2.2), Inches(1.5), Inches(0.5), title, 9, True, BLACK, PP_ALIGN.CENTER)
    txt(s, Inches(x), Inches(2.7), Inches(1.5), Inches(1.4), desc, 8, False, GRAY, PP_ALIGN.CENTER)
    if x < 10:
        arrow(s, Inches(x + 1.5), Inches(2.8), Inches(0.12), Inches(0.15))
    x += 1.62

txt(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.8), "Total time from commit to compliance assessment: < 1 second\nAll autonomous — no human intervention required", 14, True, GREEN)

# ===================== SLIDE 8: Market Gap =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Why This Does Not Exist", 32, True, BLACK)
txt(s, Inches(0.8), Inches(1.0), Inches(11), Inches(0.4), "Verified against 30+ products. No tool does cross-domain compliance reasoning from delivery events.", 14, False, GRAY)

gaps = [
    ("Compliance tools (Vanta, Drata)", "Evidence collection for one company", "No cross-domain reasoning"),
    ("GRC platforms (ServiceNow)", "Risk registers, control tracking", "Static, no event reaction"),
    ("Security tools (Snyk, Checkmarx, Wiz)", "Detect vulnerabilities", "No business impact reasoning"),
    ("SOC platforms (Panther, CrowdStrike)", "Incident response automation", "Security only, no SOX/regulatory"),
    ("Audit tools (Fieldguide, AuditBoard)", "Auditor workflow management", "Don't generate evidence"),
    ("Regulatory tools (Regology)", "Track regulatory changes", "Don't simulate per-client impact"),
    ("AI models (GPT, Claude, Copilot)", "Answer questions when asked", "Don't sit on ecosystem, don't react"),
]

txt(s, Inches(0.8), Inches(1.5), Inches(3.8), Inches(0.3), "Category", 11, True, BLUE)
txt(s, Inches(4.6), Inches(1.5), Inches(3.5), Inches(0.3), "What They Do", 11, True, BLUE)
txt(s, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.3), "What They Miss", 11, True, BLUE)

y = 1.9
for cat, does, misses in gaps:
    txt(s, Inches(0.8), Inches(y), Inches(3.8), Inches(0.4), cat, 11, False, BLACK)
    txt(s, Inches(4.6), Inches(y), Inches(3.5), Inches(0.4), does, 11, False, GRAY)
    txt(s, Inches(8.3), Inches(y), Inches(4.5), Inches(0.4), misses, 11, False, RED)
    y += 0.5

box(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.0), LIGHT_BLUE, BLUE)
multi(s, Inches(1.0), Inches(5.8), Inches(11), Inches(0.9), [
    "The gap: No product combines ecosystem monitoring + cross-domain policy evaluation +",
    "autonomous evidence generation + per-client context + deployment blocking in one agent."
], 13, BLACK)


# ===================== SLIDE 9: Business Value =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)
txt(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), "Business Value", 32, True, BLACK)

metrics = [("70%", "Audit prep\ntime reduced"), ("$2M+", "Penalties\navoided/year"), ("80%", "Findings caught\nbefore audit"), ("175%", "Year 1 ROI"), ("<1s", "Event to\nassessment")]
x = 0.5
for val, label in metrics:
    box(s, Inches(x), Inches(1.2), Inches(2.3), Inches(1.6), LIGHT_BLUE, BLUE)
    txt(s, Inches(x), Inches(1.3), Inches(2.3), Inches(0.8), val, 30, True, BLUE, PP_ALIGN.CENTER)
    txt(s, Inches(x), Inches(2.1), Inches(2.3), Inches(0.6), label, 12, False, GRAY, PP_ALIGN.CENTER)
    x += 2.55

multi(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5), [
    "Direct Savings:",
    "   Audit preparation reduced from weeks to always-ready",
    "   SLA breach penalties prevented ($50K per incident avoided)",
    "   Developer productivity gained (20-30% less compliance overhead)",
    "",
    "Revenue Opportunity:",
    "   Offer as premium managed service to clients ($20K-$50K/month)",
    "   Competitive differentiator: no competitor in IT services has this",
    "   Platform can be licensed to other services companies (SaaS model)",
    "",
    "Strategic:",
    "   Works across any vertical (mortgage, healthcare, retail, manufacturing)",
    "   Builds on existing Bedrock/Anthropic partnership",
    "   First mover in Delivery Compliance Intelligence category",
], 13, GRAY)

# ===================== SLIDE 10: Closing =====================
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
box(s, Inches(0), Inches(0), Inches(13.333), Inches(0.12), BLUE)

txt(s, Inches(2), Inches(2.2), Inches(9), Inches(0.8), "KAVACH AI", 48, True, BLACK, PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(3.2), Inches(9), Inches(0.6), "Kontinuous Audit & Vulnerability Analysis for Compliant Hardening", 22, False, GRAY, PP_ALIGN.CENTER)

box(s, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.5), LIGHT_BLUE, BLUE)
multi(s, Inches(2.8), Inches(4.4), Inches(7.8), Inches(1.3), [
    "An autonomous agent that monitors the delivery ecosystem,",
    "triggers cross-domain compliance checks on every change,",
    "and generates audit evidence — all in under one second.",
    "Domain-agnostic. Works in any vertical. No tool like this exists today.",
], 14, BLACK)

txt(s, Inches(2), Inches(6.0), Inches(9), Inches(0.5), "Working prototype available for live demo.", 15, True, GREEN, PP_ALIGN.CENTER)
txt(s, Inches(2), Inches(6.5), Inches(9), Inches(0.4), "One commit. Six domains. Zero breaches.", 16, False, BLUE, PP_ALIGN.CENTER)


# ===================== SAVE =====================
import os
path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KAVACH_AI_Clean.pptx")
prs.save(path)
print(f"Saved: {path}")
print(f"Slides: {len(prs.slides)}")
