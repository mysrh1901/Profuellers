#!/usr/bin/env python3
"""
Create KAVACH AI presentation using official PPT Template.
Uses proper layouts, placeholders, and positioning from the template.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation('/Users/Research/Downloads/Profuellers/PPT Template.pptx')

# Remove the existing sample slide
rId = prs.slides._sldIdLst[0].rId
prs.part.drop_rel(rId)
del prs.slides._sldIdLst[0]

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x5A, 0x5A, 0x5A)
BLUE = RGBColor(0x00, 0x47, 0xAB)


def add_content_slide(title_text, subtitle_text, body_lines):
    """Add a slide using 'Blank Layout with Sub Title' (Layout 2)."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    # Layout 2: PH idx=15 at L=0.7 T=0.5 (title), PH idx=0 at L=0.7 T=1.0 (subtitle)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 15:
            ph.text = title_text
            for para in ph.text_frame.paragraphs:
                para.font.size = Pt(22)
                para.font.bold = True
                para.font.color.rgb = BLACK
        elif ph.placeholder_format.idx == 0:
            ph.text = subtitle_text
            for para in ph.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.color.rgb = GRAY

    # Add body content below
    if body_lines:
        tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.font.name = "Calibri"
            p.space_after = Pt(5)
    return slide


# ================================================================
# SLIDE 1: Cover
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[0])  # Cover - Light Background
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "KAVACH AI"
        ph.text_frame.paragraphs[0].font.size = Pt(36)
        ph.text_frame.paragraphs[0].font.bold = True

# Add tagline manually
tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.2), Inches(5.6), Inches(0.5))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Knowledge-driven Audit, Vulnerability Analysis & Compliance Health"
p.font.size = Pt(16)
p.font.color.rgb = GRAY

tb = slide.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(5.6), Inches(0.4))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "One commit. Six domains. Zero breaches."
p.font.size = Pt(12)
p.font.color.rgb = BLUE

tb = slide.shapes.add_textbox(Inches(0.9), Inches(6.0), Inches(5.6), Inches(0.3))
tf = tb.text_frame
p = tf.paragraphs[0]
p.text = "Agentic Arena 2026"
p.font.size = Pt(11)
p.font.color.rgb = GRAY


# ================================================================
# SLIDE 2: Problem Statement
# ================================================================
add_content_slide(
    "Problem Statement",
    "Silent multi-domain compliance violations in regulated delivery",
    [
        "A single code commit can simultaneously violate multiple compliance domains:",
        "",
        "     SOX                    Unapproved change to a financial system",
        "     Application Security   Unresolved vulnerability deployed to production",
        "     Regulatory             Incorrect rate/APR calculation (TILA, RESPA)",
        "     Contractual SLA        48-hour remediation window missed ($50K penalty)",
        "     Fair Lending            Pricing logic changed without impact testing (ECOA)",
        "     Privacy                 PII crosses geographic boundary (GDPR, CCPA)",
        "",
        "These are checked by different tools, owned by different teams, on different cycles.",
        "Nobody connects the causal chain between them.",
        "",
        "Violations discovered months later during audits.",
        "Cost per incident: $50K - $500K.",
        "Annual risk exposure: $5M - $10M per organization.",
    ]
)

# ================================================================
# SLIDE 3: Solution Overview
# ================================================================
add_content_slide(
    "Solution: KAVACH AI",
    "An autonomous agent that monitors the ecosystem and triggers all compliance checks",
    [
        "An agent that sits on the delivery ecosystem (Git, Jenkins, Jira, Docker, AWS)",
        "and when any change happens, evaluates ALL compliance domains simultaneously.",
        "",
        "It does not replace existing tools. It connects them into one intelligent layer.",
        "",
        "",
        "Processing Flow:",
        "",
        "     Event Occurs               Message Queue              Agent Layer                    Output",
        "     ─────────────              ──────────────             ────────────                  ──────",
        "     Code commit                 ActiveMQ / SQS             Chain Reactor evaluates       Block or Allow",
        "     Infra change           ►    JMS Queue              ►   all 8 compliance         ►   Audit evidence",
        "     Access change               compliance-events          domains in < 1 second        Score update",
        "     Regulatory update                                                                    Alert teams",
        "",
        "",
        "All 8 compliance domains evaluated per event. Fully autonomous. No human intervention.",
    ]
)

# ================================================================
# SLIDE 4: 8 Compliance Domains
# ================================================================
add_content_slide(
    "8 Compliance Domains Checked Per Event",
    "All triggered simultaneously from a single change — no silo, no gap",
    [
        "",
        "     #    Domain                       What It Checks",
        "     ─    ──────                       ──────────────",
        "     1    SOX ITGC                     Dual approval, segregation of duties, change docs",
        "     2    Application Security         SAST/DAST findings, CVE tracking, deployment gate",
        "     3    Regulatory                   TILA rate accuracy, RESPA, GDPR residency, DORA",
        "     4    Fair Lending                 Disparate impact when pricing/eligibility changes",
        "     5    Contractual (MSA/SLA)        Notification deadlines, SLA timers, penalty triggers",
        "     6    Privacy                      PII detection, consent, cross-border transfers",
        "     7    Infrastructure               Cloud misconfigs, IAM drift, encryption, backups",
        "     8    Audit Evidence               Auto-generated narrative for every change",
        "",
        "",
        "Policy-driven: adding a new check = adding a policy object.",
        "No code changes to the agents required.",
    ]
)


# ================================================================
# SLIDE 5: Architecture
# ================================================================
add_content_slide(
    "Architecture",
    "Event-driven, policy-based, domain-agnostic agent pipeline",
    [
        "",
        "     ┌──────────────┐      ┌───────────────┐      ┌──────────────────────┐      ┌─────────────────┐",
        "     │              │      │               │      │                      │      │                 │",
        "     │  EVENT       │      │  MESSAGE      │      │  AGENT LAYER         │      │  OUTPUT         │",
        "     │  SOURCES     │─────►│  QUEUE        │─────►│                      │─────►│                 │",
        "     │              │      │               │      │  Chain Reactor Agent  │      │  Block / Allow  │",
        "     │  Git Hooks   │      │  ActiveMQ     │      │    Evaluates policies │      │  Audit Evidence │",
        "     │  Jenkins     │      │  (JMS)        │      │    across 8 domains   │      │  Score Update   │",
        "     │  Jira        │      │               │      │                      │      │  Dashboard      │",
        "     │  Docker      │      │               │      │  Audit Narrator Agent │      │  Alerts         │",
        "     │  AWS Config  │      │               │      │    Generates evidence │      │                 │",
        "     │  API / Postman│      │               │      │    (LLM optional)    │      │                 │",
        "     │              │      │               │      │                      │      │                 │",
        "     └──────────────┘      └───────────────┘      │  Policy Engine        │      └─────────────────┘",
        "                                                   │    15 rules            │",
        "                                                   │    9 verticals         │",
        "                                                   └──────────────────────┘",
        "",
        "     Data Layer: Database (H2 / PostgreSQL) — Events, Findings, Narratives, Client Profiles",
        "     LLM Layer: AWS Bedrock (Claude) — optional, for code understanding and narrative quality",
    ]
)

# ================================================================
# SLIDE 6: Works Across Verticals
# ================================================================
add_content_slide(
    "Works Across Any Vertical",
    "Same agents, different policies. Zero code changes for a new industry.",
    [
        "",
        "     Vertical                         Applicable Frameworks",
        "     ────────                         ─────────────────────",
        "     Mortgage / Banking               SOX, TILA, RESPA, ECOA, PCI-DSS, GLBA",
        "     Healthcare / Pharma              HIPAA, HITECH, FDA 21 CFR Part 11",
        "     Insurance                        SOX, NAIC Model Laws, CCPA, State Regs",
        "     Retail / E-commerce              PCI-DSS, CCPA, GDPR, SOX",
        "     Manufacturing                    ISO 27001, NIST, ITAR, Export Controls",
        "     Public Sector                    FedRAMP, FISMA, CMMC",
        "",
        "",
        "Currently: 15 policies loaded across 9 verticals and 11 compliance domains.",
        "",
        "To support a new client in a new industry:",
        "     Add policies to the Policy Engine configuration.",
        "     No changes to the Chain Reactor Agent or Audit Narrator Agent.",
        "     Deploy. Done.",
    ]
)

# ================================================================
# SLIDE 7: Market Gap
# ================================================================
add_content_slide(
    "This Does Not Exist in the Market",
    "Verified against 30+ products and platforms",
    [
        "",
        "     Category                        What They Do                    What They Miss",
        "     ────────                        ────────────                    ──────────────",
        "     Compliance (Vanta, Drata)        Evidence for one company        No cross-domain reasoning",
        "     GRC (ServiceNow, Archer)         Risk registers                  Static, no event reaction",
        "     Security (Snyk, Checkmarx)       Detect vulnerabilities          No business impact",
        "     SOC (Panther, CrowdStrike)       Incident response               Security only, no SOX/reg",
        "     Audit (Fieldguide, AuditBoard)   Auditor workflow                Don't generate evidence",
        "     Regulatory (Regology)            Track reg changes               No per-client simulation",
        "     AI Models (GPT, Claude)          Answer when asked               Don't react autonomously",
        "",
        "",
        "The gap nobody fills:",
        "",
        "     Ecosystem monitoring  +  Cross-domain policy evaluation  +",
        "     Autonomous evidence generation  +  Per-client context  +  Deployment blocking",
        "",
        "     No single product combines these capabilities today.",
    ]
)


# ================================================================
# SLIDE 8: Business Value
# ================================================================
add_content_slide(
    "Business Value",
    "Quantified impact across savings, revenue, and strategic positioning",
    [
        "",
        "     Metric                                    Impact",
        "     ──────                                    ──────",
        "     Audit preparation time                    Reduced by 70%",
        "     SLA breach penalties avoided              $500K - $2M per year",
        "     Findings caught before audit              80%",
        "     Developer compliance overhead             Reduced 20-30%",
        "     Event to compliance assessment            Under 1 second",
        "     Year 1 ROI (conservative)                175%",
        "",
        "",
        "     Revenue Opportunity:",
        "         Premium managed service to clients: $20K - $50K / month per engagement",
        "         Platform licensing to other IT services companies",
        "         First mover in Delivery Compliance Intelligence category",
        "",
        "     Strategic Advantage:",
        "         Domain-agnostic — works in any regulated industry",
        "         No competitor in IT services has this capability today",
    ]
)

# ================================================================
# SLIDE 9: Section Divider - Demo
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[7])  # Section Divider
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Live Demo"
        ph.text_frame.paragraphs[0].font.size = Pt(32)
        ph.text_frame.paragraphs[0].font.bold = True
    elif ph.placeholder_format.idx == 15:
        ph.text = "Working prototype available"
        ph.text_frame.paragraphs[0].font.size = Pt(14)

# ================================================================
# SLIDE 10: Thank You
# ================================================================
slide = prs.slides.add_slide(prs.slide_layouts[11])  # Thank You layout
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "KAVACH AI"
        ph.text_frame.paragraphs[0].font.size = Pt(28)
        ph.text_frame.paragraphs[0].font.bold = True

tb = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(5.5), Inches(1.5))
tf = tb.text_frame
tf.word_wrap = True
lines = [
    "Knowledge-driven Audit, Vulnerability Analysis & Compliance Health",
    "",
    "One commit. Six domains. Zero breaches.",
    "",
    "Working prototype ready for demo.",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(13)
    p.font.color.rgb = GRAY
    p.font.name = "Calibri"


# ================================================================
# SAVE
# ================================================================
import os
path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KAVACH_AI_Arena.pptx")
prs.save(path)
print(f"Saved: {path}")
print(f"Slides: {len(prs.slides)}")
